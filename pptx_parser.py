import pandas as pd
import yaml
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE
from math import sqrt
from typing import Dict, List, Any, Tuple
from text_utils import extract_details_from_title

def emu_to_cm(emu: float) -> float:
    return round(emu / 360000.0, 2)

def get_shape_layout(shape) -> Dict[str, float]:
    return {
        "x": emu_to_cm(shape.left),
        "y": emu_to_cm(shape.top),
        "width": emu_to_cm(shape.width),
        "height": emu_to_cm(shape.height),
    }

def get_shape_center(shape) -> Tuple[float, float]:
    return shape.left + shape.width / 2, shape.top + shape.height / 2

def table_shape_to_df(shape) -> "pd.DataFrame | None":
    if not hasattr(shape, "table"):
        return None
    table = shape.table
    rows = table.rows
    cols = table.columns

    data = []
    for r in range(len(rows)):
        row_vals = []
        for c in range(len(cols)):
            cell = table.cell(r, c)
            txt = cell.text_frame.text if cell.text_frame else ""
            row_vals.append(txt.strip())
        data.append(row_vals)

    if data:
        header = data[0]
        body = data[1:] if len(data) > 1 else []
        if len(set(h or f"col_{i}" for i, h in enumerate(header))) == len(header) and any(h.strip() for h in header):
            df = pd.DataFrame(body, columns=[h if h else f"col_{i}" for i, h in enumerate(header)])
        else:
            df = pd.DataFrame(data)
    else:
        df = pd.DataFrame()
    return df


def chart_shape_to_df(shape) -> "pd.DataFrame | None":
    if not hasattr(shape, "chart"):
        return None
    chart = shape.chart

    categories = []
    try:
        if chart.plots and chart.plots[0].categories is not None:
            for cat in chart.plots[0].categories:
                if hasattr(cat, "label"):
                    categories.append(str(cat.label))
                else:
                    categories.append(str(cat))
        else:
            categories = None
    except Exception:
        categories = None

    series_data = {}
    max_len = 0
    for s in chart.series:
        name = s.name if s.name is not None else f"series_{len(series_data)}"
        values = []
        for v in (s.values or []):
            if hasattr(v, "value"):
                values.append(v.value)
            else:
                try:
                    values.append(float(v))
                except Exception:
                    values.append(v)
        series_data[str(name)] = values
        max_len = max(max_len, len(values))

    if categories is None:
        categories = [f"cat_{i+1}" for i in range(max_len)]
    else:
        if len(categories) < max_len:
            categories = categories + [f"cat_{i+1}" for i in range(len(categories), max_len)]
        elif len(categories) > max_len:
            categories = categories[:max_len]

    df = pd.DataFrame({"category": categories})
    for series_name, vals in series_data.items():
        if len(vals) < max_len:
            vals = vals + [None] * (max_len - len(vals))
        elif len(vals) > max_len:
            vals = vals[:max_len]
        df[series_name] = vals
    return df


class PptxParser:

    def __init__(self, pptx_path: Path):
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPT file not found: {pptx_path}")
        self.presentation = Presentation(pptx_path)
        self.file_name = pptx_path.stem

    def _get_shape_type(self, shape) -> str:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if hasattr(shape, "chart"):
                chart_type = shape.chart.chart_type
                if chart_type in (
                    XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.COLUMN_STACKED_100,
                    XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.BAR_STACKED_100
                ):
                    return "chart-bar"
                elif chart_type in (
                    XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS, XL_CHART_TYPE.LINE_STACKED, 
                    XL_CHART_TYPE.LINE_STACKED_100, XL_CHART_TYPE.LINE_MARKERS_STACKED, XL_CHART_TYPE.LINE_MARKERS_STACKED_100
                ):
                    return "chart-line"
                else:
                    return "chart-other"
            return "chart"

        if shape.has_text_frame and shape.text.strip():
            return "text"
            
        return "other"

    def _classify_shapes(self, slide) -> Dict[str, List[Dict[str, Any]]]:
        classified_shapes = {
            "text": [], "table": [], "chart-bar": [], "chart-line": [], 
            "chart-other": [], "other": []
        }
        
        for shape in slide.shapes:
            shape_type = self._get_shape_type(shape)
            if shape_type == "other":
                continue

            shape_info = {
                "layout": get_shape_layout(shape),
                "center": get_shape_center(shape),
                "obj": shape
            }
            if shape_type == "text":
                shape_info["content"] = shape.text.strip()
            elif shape_type == "table":
                try:
                    shape_info["dataframe"] = table_shape_to_df(shape)
                except Exception as e:
                    shape_info["dataframe"] = None
            elif shape_type in ("chart-bar", "chart-line", "chart-other"):
                try:
                    shape_info["dataframe"] = chart_shape_to_df(shape)
                except Exception as e:
                    shape_info["dataframe"] = None

            if shape_type not in classified_shapes:
                classified_shapes[shape_type] = []
            classified_shapes[shape_type].append(shape_info)

        classified_shapes["text"].sort(key=lambda s: s["layout"]["y"])
        return classified_shapes

    def parse_slide(self, slide_idx: int = 0) -> Dict[str, Any]:
        if slide_idx >= len(self.presentation.slides):
            raise IndexError(f"Slide index {slide_idx} out of range.")
        
        slide = self.presentation.slides[slide_idx]

        shapes = self._classify_shapes(slide)
        text_shapes = shapes["text"]
        content_elements_shapes = (
            shapes["table"] + shapes["chart-bar"] + 
            shapes["chart-line"] + shapes["chart-other"]
        )

        slide_title, analysis_text, element_titles = None, None, []
        if len(text_shapes) >= 1:
            slide_title = {"content": text_shapes[0]["content"], "layout": text_shapes[0]["layout"]}
        if len(text_shapes) >= 2:
            analysis_text = {"content": text_shapes[1]["content"], "layout": text_shapes[1]["layout"]}
            element_titles = text_shapes[2:]
        elif len(text_shapes) == 1:
            element_titles = text_shapes[1:]

        content_elements = []
        for title_info in element_titles:
            if not content_elements_shapes:
                break
            
            closest_element = min(
                content_elements_shapes,
                key=lambda el: sqrt(
                    (el["center"][0] - title_info["center"][0])**2 +
                    (el["center"][1] - title_info["center"][1])**2
                )
            )
            
            details = extract_details_from_title(title_info["content"])
            
            element = {
                "title": {"content": title_info["content"], "layout": title_info["layout"]},
                "data": closest_element["dataframe"],
                "shape_type": self._get_shape_type(closest_element["obj"]),
                "layout": closest_element["layout"],
                **details
            }
            content_elements.append(element)
            content_elements_shapes.remove(closest_element)

        final_structure = {
            "template_slide": {
                "slide_size": {"width": emu_to_cm(self.presentation.slide_width), "height": emu_to_cm(self.presentation.slide_height)},
                "title": slide_title,
                "analysis": analysis_text,
                "content_elements": content_elements
            }
        }
        return final_structure

    @staticmethod
    def save_dict_as_yaml(data: Dict, output_path: Path):
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, default_flow_style=False, allow_unicode=True, indent=2, sort_keys=False)
        print(f"Successfully saved extracted structure to: {output_path}")

