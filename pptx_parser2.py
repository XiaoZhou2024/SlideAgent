# pptx_parser.py
import pandas as pd
import yaml
import os
import shutil
import json
import io
import base64

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE  # 导入图表类型枚举
from pptx.util import Cm
from math import sqrt, hypot
from typing import Dict, List, Any, Tuple
from pptxtopdf import convert
from pdf2image import convert_from_path

# 从其他模块导入辅助函数
from text_utils import extract_details_from_title
from pptx_analyser import _call_vision_model_v2


# --- 辅助函数 ---
def emu_to_cm(emu: float) -> float:
    """将EMU（英国测量单位）转换为厘米，并保留两位小数。"""
    return round(emu / 360000.0, 2)


def get_shape_layout(shape) -> Dict[str, float]:
    """从shape对象提取布局信息（位置和尺寸），单位为厘米。"""
    return {
        "x": emu_to_cm(shape.left),
        "y": emu_to_cm(shape.top),
        "width": emu_to_cm(shape.width),
        "height": emu_to_cm(shape.height),
    }


def get_shape_center(shape) -> Tuple[float, float]:
    """计算形状中心的坐标（EMU单位）。"""
    return shape.left + shape.width / 2, shape.top + shape.height / 2


def table_shape_to_df(shape) -> "pd.DataFrame | None":
    """将表格形状提取为 DataFrame。"""
    if not hasattr(shape, "table"):
        return None
    table = shape.table
    rows = table.rows
    cols = table.columns

    # 读取所有单元格文本
    data = []
    for r in range(len(rows)):
        row_vals = []
        for c in range(len(cols)):
            cell = table.cell(r, c)
            # 去掉 PPT 中常见的换行和空白
            txt = cell.text_frame.text if cell.text_frame else ""
            row_vals.append(txt.strip())
        data.append(row_vals)

    # 简单规则：如果第一行是列名，尝试做表头
    if data:
        header = data[0]
        body = data[1:] if len(data) > 1 else []
        # 检查表头是否有重复或为空，若有则退化为无表头
        if len(set(h or f"col_{i}" for i, h in enumerate(header))) == len(header) and any(h.strip() for h in header):
            df = pd.DataFrame(body, columns=[h if h else f"col_{i}" for i, h in enumerate(header)])
        else:
            df = pd.DataFrame(data)
    else:
        df = pd.DataFrame()
    return df


def chart_shape_to_df(shape) -> "pd.DataFrame | None":
    """将图表形状（柱状/条形/折线等）提取为 DataFrame，列包含类别和各系列。"""
    if not hasattr(shape, "chart"):
        return None
    chart = shape.chart

    # 读取分类轴（通常是X轴的类别）
    categories = []
    try:
        if chart.plots and chart.plots[0].categories is not None:
            for cat in chart.plots[0].categories:
                # cat 可能是字符串、数字或包含 .label 的对象
                if hasattr(cat, "label"):
                    categories.append(str(cat.label))
                else:
                    categories.append(str(cat))
        else:
            # 没有显式类别时，使用序号
            # 稍后根据系列数据长度动态填充
            categories = None
    except Exception:
        categories = None

    series_data = {}
    max_len = 0
    for s in chart.series:
        name = s.name if s.name is not None else f"series_{len(series_data)}"
        values = []
        # s.values 通常是一个序列，元素可能为 chart data point 对象或原始值
        for v in (s.values or []):
            # data point 对象可通过 .value 拿数值
            if hasattr(v, "value"):
                values.append(v.value)
            else:
                try:
                    values.append(float(v))
                except Exception:
                    values.append(v)
        series_data[str(name)] = values
        max_len = max(max_len, len(values))

    # 构造 DataFrame
    if categories is None:
        categories = [f"cat_{i + 1}" for i in range(max_len)]
    else:
        # 类别数量与最长系列长度不一致时，做长度对齐
        if len(categories) < max_len:
            categories = categories + [f"cat_{i + 1}" for i in range(len(categories), max_len)]
        elif len(categories) > max_len:
            categories = categories[:max_len]

    df = pd.DataFrame({"category": categories})
    for series_name, vals in series_data.items():
        # 补齐或截断
        if len(vals) < max_len:
            vals = vals + [None] * (max_len - len(vals))
        elif len(vals) > max_len:
            vals = vals[:max_len]
        df[series_name] = vals
    return df


def _convert_ppt_to_image(ppt_path: str, slide_number: int, output_folder_path: str):
    """将指定 PPT 文件的特定页码转换为图像。"""
    file_name = os.path.splitext(os.path.basename(ppt_path))[0]
    file_path = Path(ppt_path).parent
    temp_pdf_output_path = os.path.join(file_path, f"{file_name}_temp_pdf")
    temp_image_output_path = os.path.join(file_path, f"{file_name}_temp_images")
    os.makedirs(temp_pdf_output_path, exist_ok=True)
    os.makedirs(temp_image_output_path, exist_ok=True)

    try:
        convert(input_path=ppt_path, output_folder_path=temp_pdf_output_path)
        pdf_file_path = os.path.join(temp_pdf_output_path, f"{file_name}.pdf")
        if not os.path.exists(pdf_file_path):
            raise FileNotFoundError(f"PDF 文件转换失败: {pdf_file_path}")
        # images = convert_from_path(pdf_file_path)
        images = convert_from_path(
            pdf_file_path,
            poppler_path=r"C:\poppler-25.07.0\Library\bin",  # 改成你的实际路径
            fmt="jpeg",  # 或 png
            dpi=200
        )
        if not (0 <= slide_number - 1 < len(images)):
            print(f"警告: 页码 {slide_number} 超出范围。")
            return None, None, None, None
        target_slide_pil_image = images[slide_number - 1]
        width_px, height_px = target_slide_pil_image.size
        return target_slide_pil_image, width_px, height_px, temp_image_output_path
    except Exception as e:
        print(f"PPT 到图像转换失败: {e}")
        return None, None, None, None
    finally:
        if os.path.exists(temp_pdf_output_path):
            shutil.rmtree(temp_pdf_output_path)


# --- 主解析器类 ---
class PptxParser:
    """
    一个用于解析PPTX文件幻灯片并提取其结构化信息的类。
    """

    def __init__(self, pptx_path: Path):
        """
        初始化解析器。

        Args:
            pptx_path: 要解析的PPTX文件的路径。
        """
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPT文件未找到: {pptx_path}")
        self.presentation = Presentation(pptx_path)
        self.file_name = pptx_path.stem
        self.file_path = str(pptx_path)
        self.slide_count = len(self.presentation.slides)

    def _get_shape_type(self, shape) -> str:
        """
        获取形状的规范化类型名称。
        """
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"

        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            if hasattr(shape, "chart"):
                chart_type = shape.chart.chart_type
                # 柱状图和条形图
                if chart_type in (
                        XL_CHART_TYPE.COLUMN_CLUSTERED, XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.COLUMN_STACKED_100,
                        XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.BAR_STACKED_100
                ):
                    return "chart-bar"
                # 折线图
                elif chart_type in (
                        XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS, XL_CHART_TYPE.LINE_STACKED,
                        XL_CHART_TYPE.LINE_STACKED_100, XL_CHART_TYPE.LINE_MARKERS_STACKED,
                        XL_CHART_TYPE.LINE_MARKERS_STACKED_100
                ):
                    return "chart-line"
                else:
                    # 其他图表类型的回退选项
                    return "chart-other"
            return "chart"

        if shape.has_text_frame and shape.text.strip():
            return "text"

        return "other"

    def _extract_pptx_elements(self, slide) -> Dict[str, Any]:
        """
        从单个幻灯片中提取表格和图表数据及pptx坐标。
        """
        content_elements = []
        for idx, shape in enumerate(slide.shapes):
            shape_type = self._get_shape_type(shape)
            if shape_type == "text":
                content_elements.append({
                    "id": idx,
                    "type": 'textBox',
                    "text": shape.text.strip(),
                    'role': '',
                    "layout": get_shape_layout(shape)
                })
            if shape_type == "table":
                df = table_shape_to_df(shape)
                if df is not None and not df.empty:
                    content_elements.append({
                        "id": idx,
                        "type": 'table',
                        'role': '',
                        "layout": get_shape_layout(shape),
                        "data": df.to_dict(orient='records'),
                    })
            elif "chart" in shape_type:
                df = chart_shape_to_df(shape)
                if df is not None and not df.empty:
                    content_elements.append({
                        "id": idx,
                        "type": 'chart',
                        'role': '',
                        "layout": get_shape_layout(shape),
                        "data": df.to_dict(orient='records'),
                    })
        template_slide = {
                "slide_size": {"width": emu_to_cm(self.presentation.slide_width),
                               "height": emu_to_cm(self.presentation.slide_height)},
                "elements": content_elements
        }
        return template_slide

    def _get_vlm_analysis(self, slide_idx: int) -> Tuple[List[Dict[str, Any]], int, int]:
        """
        将幻灯片转换为图像并调用 VLM 获取文本和位置信息。
        返回 VLM 结果、图像宽度和高度。
        """
        output_dir = Path(".").resolve()
        img, img_w, img_h, temp_path = _convert_ppt_to_image(
            self.file_path, slide_idx + 1, str(output_dir)
        )
        if not img:
            return [], 0, 0

        buffered = io.BytesIO()
        img.save(buffered, format="JPEG")
        base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

        vlm_results = _call_vision_model_v2(base64_image)
        if temp_path and os.path.exists(temp_path):
            shutil.rmtree(temp_path)

        return vlm_results, img_w, img_h

    def _bbox_px_2_layout(self, bbox: list, pptx_cm_to_px_w_ratio: float, pptx_cm_to_px_h_ratio: float) -> dict:
        return {
            "x": round(bbox[0] / pptx_cm_to_px_w_ratio, 2),
            "y": round(bbox[1] / pptx_cm_to_px_h_ratio, 2),
            "width": round((bbox[2] - bbox[0]) / pptx_cm_to_px_w_ratio, 2),
            "height": round((bbox[3] - bbox[1]) / pptx_cm_to_px_h_ratio, 2),
        }

    def _nearest_point(self, point, points):
        """
        point: (x, y)
        points: [(x1, y1), (x2, y2), ...]
        return: (最近点, 距离, 索引)
        """
        px, py = point
        best = None
        best_dist = float('inf')
        best_idx = -1

        for i, (x, y) in enumerate(points):
            # 欧氏距离
            d = hypot(x - px, y - py)
            if d < best_dist:
                best_dist = d
                best = (x, y)
                best_idx = i
        return  best_idx

    def _match_elements(self, pptx_elements: Dict[str, Any], vlm_results: List[Dict[str, Any]], img_w: int,
                        img_h: int) -> Dict[str, Any]:
        """
        将 pptx 提取的元素与 VLM 识别的文本进行匹配和整合。
        """
        pptx_cm_to_px_w_ratio = img_w / emu_to_cm(self.presentation.slide_width)
        pptx_cm_to_px_h_ratio = img_h / emu_to_cm(self.presentation.slide_height)
        elements = pptx_elements.get('elements', [])
        points = []
        for element in elements:
            points.append((element.get('layout').get('x'), element.get('layout').get('y')))

        ### 遍历 VLM 识别的元素
        for item in vlm_results:
            item_point = (round(item["bbox"][0] / pptx_cm_to_px_w_ratio, 2), round(item["bbox"][1] / pptx_cm_to_px_h_ratio, 2))
            nearest_point_idx = self._nearest_point(item_point, points)
            elements[nearest_point_idx]['role'] = item['shape_type']

        template_slide = {
                "slide_size": pptx_elements.get('slide_size'),
                "elements": elements
        }

        return template_slide

    def parse_slide_vlm(self, slide_idx: int = 0) -> Dict[str, Any]:
        """
        使用 VLM 解析幻灯片结构，并与 pptx 数据进行匹配。
        """
        if slide_idx >= len(self.presentation.slides):
            raise IndexError(f"幻灯片索引 {slide_idx} 超出范围。")

        slide = self.presentation.slides[slide_idx]

        # 步骤 1: 从 pptx 提取元素
        pptx_elements = self._extract_pptx_elements(slide)
        print(f"从幻灯片提取元素: {pptx_elements}")

        # 步骤 2: 调用 VLM 进行分析
        vlm_results, img_w, img_h = self._get_vlm_analysis(slide_idx)
        print(f"VLM 分析结果: {vlm_results}")
        if not vlm_results:
            return {"error": "无法生成幻灯片图片或调用VLM"}

        # 步骤 3: 匹配并整合数据
        structured_data = self._match_elements(pptx_elements, vlm_results, img_w, img_h)

        template_slide = {"template_slide": structured_data}
        return template_slide

    @staticmethod
    def save_dict_as_yaml(data: Dict, output_path: Path):
        """将字典保存为YAML文件。"""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            yaml.dump(data, f, default_flow_style=False, allow_unicode=True, indent=2, sort_keys=False)
        print(f"成功将提取的结构保存到: {output_path}")


# --- 使用示例 ---
if __name__ == "__main__":
    ppt_file_path = "./ReSlide/ReSlide_01/template-1/temp/北京市良乡.pptx"  # 替换为你的PPT文件路径
    output_directory = "./output/"

    if not os.path.exists(ppt_file_path):
        print(f"错误：测试文件 {ppt_file_path} 不存在。请修改为正确的路径。")
    else:
        parser = PptxParser(Path(ppt_file_path))
        slide_to_analyze = 0  # 幻灯片索引，从0开始

        # 运行 VLM 解析并匹配
        try:
            structured_data = parser.parse_slide_vlm(slide_to_analyze)
            print("--- 幻灯片结构化解析结果 ---")
            print(json.dumps(structured_data, indent=2, ensure_ascii=False))

            # 您也可以选择保存为 YAML 文件
            parser.save_dict_as_yaml(structured_data, Path(output_directory, "slide_analysis.yaml"))
        except Exception as e:
            print(f"解析过程中发生错误: {e}")
